
import logging
import sys
import io  # 导入 io 模块
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
import docx
import pypandoc
from pptx import Presentation
import os
sys.stdout.reconfigure(encoding='utf-8')
def pdf_text(pdf_path):
    with open(pdf_path, 'rb') as file:
        rsrcmgr = PDFResourceManager()
        laparams = LAParams()
        retstr = io.StringIO()  # 使用 StringIO 存储转换的文本
        device = TextConverter(rsrcmgr, retstr, laparams=laparams)  # 去掉 codec 参数
        interpreter = PDFPageInterpreter(rsrcmgr, device)

        for page in PDFPage.get_pages(file):
            interpreter.process_page(page)

        text = retstr.getvalue()  # 获取转换后的文本
        return text


def docx_text(file_path):
    try:
        # 打开 docx 文件
        document = docx.Document(file_path)

        # 创建一个集合用于存储已提取的文本数据，以避免重复
        seen = set()
        text_data = []

        # 遍历文档中的所有段落
        for paragraph in document.paragraphs:
            if paragraph.text and paragraph.text not in seen:
                text_data.append(paragraph.text)
                seen.add(paragraph.text)

        # 遍历文档中的所有表格
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text and cell.text not in seen:
                        text_data.append(cell.text)
                        seen.add(cell.text)

        # 合并所有段落和表格的文本并返回
        return "\n".join(text_data)

    except Exception as e:
        logging.error(f"读取 DOCX 文件时出错: {str(e)}")
        return None


def txt_text(file_path):
    try:
        # 读取文本内容
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()

        # 返回文件中的每行文本
        return text
    except Exception as e:
        logging.error(f"txt文件读取错误: {str(e)}")
        return None


def ppt_text(file_path):
    try:
        # 读取文本内容
        presentation = Presentation(file_path)
        text_data = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_data.append(shape.text)

        # 返回提取的文本
        return "\n".join(text_data)
    except Exception as e:
        logging.error(f"ppt文件读取错误: {str(e)}")
        return None


def convert_doc_to_docx(input_file_path):
    try:
        # 检查输入文件是否存在
        if not os.path.exists(input_file_path):
            raise FileNotFoundError(f"输入文件 {input_file_path} 未找到")

        # 使用pypandoc将.doc文件转换为.docx并提取文本
        output_text = pypandoc.convert_file(input_file_path, 'docx',)

        # 返回提取的文本内容
        return output_text

    except Exception as e:
        logging.error(f"转换 .doc 到 .docx 时出错: {str(e)}")
        return None
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python script_name.py <文件路径>")
        sys.exit(1)
    # 从命令行参数获取文件路径
file_path = sys.argv[1]

if not os.path.exists(file_path):
    print(f"文件 {file_path} 不存在")
    sys.exit(1)

# 根据文件扩展名选择对应的方法
ext = os.path.splitext(file_path)[1].lower()
if ext == ".pdf":
    result = pdf_text(file_path)
elif ext == ".docx":
    result = docx_text(file_path)
elif ext == ".txt":
    result = txt_text(file_path)
elif ext in [".ppt", ".pptx"]:
    result = ppt_text(file_path)
elif ext == ".doc":
    result = convert_doc_to_docx(file_path)
else:
    print(f"不支持的文件类型: {ext}")
    sys.exit(1)

# 输出结果
if result:
    os.remove(file_path)
    print(result)
else:
    print("未提取到任何内容，或发生错误")